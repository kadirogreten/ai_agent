#!/usr/bin/env python3
"""
AgentArmy — Kimya Sektörü müşteri sunumu (PPTX) üretici.
sales/build_deck.py ile aynı marka stili (açık zemin, teal + turuncu).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- Marka renk paleti -------------------------------------------------
BG_DARK   = RGBColor(0xF3, 0xF6, 0xFA)   # ferah açık zemin
BG_PANEL  = RGBColor(0xFF, 0xFF, 0xFF)   # beyaz kartlar
ACCENT    = RGBColor(0x0E, 0x9C, 0x96)   # teal (birincil vurgu)
ACCENT2   = RGBColor(0x2F, 0xA8, 0x55)   # yeşil
ACCENT3   = RGBColor(0xE2, 0x73, 0x1F)   # sıcak turuncu
WHITE     = RGBColor(0x16, 0x22, 0x32)   # ana koyu metin
MUTED     = RGBColor(0x57, 0x64, 0x78)   # ikincil gri metin
CARD_LINE = RGBColor(0xDD, 0xE4, 0xEF)   # açık kart kenarlığı

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT = "Calibri"


def fix(text):
    try:
        text = text.encode("utf-16", "surrogatepass").decode("utf-16")
    except Exception:
        pass
    return "".join(c for c in text if not (0xD800 <= ord(c) <= 0xDFFF))


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_DARK
    bg.line.fill.background()
    bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element)
    s.shapes._spTree.insert(2, bg._element)
    return s


def txt(slide, left, top, width, height, text, size=18, color=WHITE,
        bold=False, align=PP_ALIGN.LEFT, font=FONT, anchor=MSO_ANCHOR.TOP,
        line_spacing=1.0, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = fix(text).split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
        r.font.color.rgb = color
    return tb


def bullets(slide, left, top, width, height, items, size=18, color=WHITE,
            gap=6, bullet_color=ACCENT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        rb = p.add_run()
        rb.text = "▸  "
        rb.font.size = Pt(size)
        rb.font.color.rgb = bullet_color
        rb.font.name = FONT
        rb.font.bold = True
        if isinstance(it, tuple):
            lead, rest = it
            r1 = p.add_run(); r1.text = fix(lead)
            r1.font.size = Pt(size); r1.font.bold = True
            r1.font.color.rgb = color; r1.font.name = FONT
            r2 = p.add_run(); r2.text = fix(rest)
            r2.font.size = Pt(size); r2.font.color.rgb = MUTED
            r2.font.name = FONT
        else:
            r = p.add_run(); r.text = fix(it)
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.name = FONT
    return tb


def panel(slide, left, top, width, height, fill=BG_PANEL, line=CARD_LINE,
          radius=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def accent_bar(slide, left, top, width=Inches(0.9), color=ACCENT):
    b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Pt(6))
    b.fill.solid(); b.fill.fore_color.rgb = color
    b.line.fill.background(); b.shadow.inherit = False
    return b


def kicker(slide, text, color=ACCENT):
    txt(slide, Inches(0.7), Inches(0.45), Inches(10), Inches(0.4),
        text.upper(), size=13, color=color, bold=True)


def title(slide, text, size=30):
    txt(slide, Inches(0.7), Inches(0.8), Inches(12), Inches(1.0),
        text, size=size, color=WHITE, bold=True)
    accent_bar(slide, Inches(0.72), Inches(1.62))


def card(slide, left, top, width, height, icon, head, body,
         accent=ACCENT, head_size=16, body_size=12.5):
    p = panel(slide, left, top, width, height)
    txt(slide, left + Inches(0.25), top + Inches(0.18), width - Inches(0.5),
        Inches(0.5), icon, size=22, color=accent, bold=True)
    txt(slide, left + Inches(0.25), top + Inches(0.62), width - Inches(0.5),
        Inches(0.5), head, size=head_size, color=WHITE, bold=True)
    txt(slide, left + Inches(0.25), top + Inches(1.05), width - Inches(0.5),
        height - Inches(1.15), body, size=body_size, color=MUTED,
        line_spacing=1.08)
    return p


def step_flow(slide, items, yy, cw=Inches(2.18), ch=Inches(2.1),
              cgap=Inches(0.18), color=ACCENT2, head_size=15, body_size=11.5):
    n = len(items)
    totw = n * int(cw) + (n - 1) * int(cgap)
    startx = int((int(SLIDE_W) - totw) / 2)
    for i, (ic, hd, bd) in enumerate(items):
        cx = Emu(startx + i * (int(cw) + int(cgap)))
        panel(slide, cx, yy, cw, ch, fill=BG_PANEL)
        txt(slide, cx, yy + Inches(0.22), cw, Inches(0.6), ic, size=26,
            color=color, align=PP_ALIGN.CENTER)
        txt(slide, cx, yy + Inches(0.82), cw, Inches(0.5), hd, size=head_size,
            color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        txt(slide, Emu(int(cx) + int(Inches(0.12))), yy + Inches(1.28),
            Emu(int(cw) - int(Inches(0.24))), ch - Inches(1.35), bd,
            size=body_size, color=MUTED, align=PP_ALIGN.CENTER,
            line_spacing=1.05)
        if i < n - 1:
            ar = slide.shapes.add_shape(
                MSO_SHAPE.CHEVRON,
                Emu(int(cx) + int(cw) - int(Inches(0.02))),
                yy + Inches(0.8), Inches(0.22), Inches(0.5))
            ar.fill.solid(); ar.fill.fore_color.rgb = ACCENT
            ar.line.fill.background(); ar.shadow.inherit = False


# =======================================================================
# SLIDE 1 — KAPAK
# =======================================================================
s = add_slide()
strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.25), SLIDE_H)
strip.fill.solid(); strip.fill.fore_color.rgb = ACCENT
strip.line.fill.background(); strip.shadow.inherit = False

txt(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(0.5),
    "AGENTARMY  ·  KİMYA SEKTÖRÜ ÇÖZÜMLERİ", size=15,
    color=ACCENT, bold=True)
txt(s, Inches(0.9), Inches(1.8), Inches(11.8), Inches(2.0),
    "Kimya Sektörü İçin\nYapay Zeka Ajan Ordusu", size=46,
    color=WHITE, bold=True, line_spacing=1.0)
txt(s, Inches(0.9), Inches(3.9), Inches(11.0), Inches(1.3),
    "Pazar istihbaratından tedarik otomasyonuna, regülasyon takibinden\n"
    "Ar-Ge dokümantasyonuna — uzman ajan ekipleri sizin yerinize çalışır.\n"
    "Kanıt, izlenebilirlik ve insan onayı varsayılan olarak gelir.",
    size=17, color=MUTED, line_spacing=1.25)

badges = ["Üretim & Pazar İstihbaratı", "Tedarik Otomasyonu",
          "Regülasyon & Uyum Takibi", "Ar-Ge & Teknik Doküman"]
bx = Inches(0.9)
for b in badges:
    w = Inches(2.95)
    panel(s, bx, Inches(5.55), w, Inches(0.6), fill=BG_PANEL)
    txt(s, bx, Inches(5.68), w, Inches(0.4), b, size=12, color=WHITE,
        bold=True, align=PP_ALIGN.CENTER)
    bx = Emu(int(bx) + int(w) + int(Inches(0.12)))

txt(s, Inches(0.9), Inches(6.7), Inches(11), Inches(0.4),
    "Müşteri Sunumu  ·  Gizli ve Özeldir", size=12, color=MUTED)

# =======================================================================
# SLIDE 2 — KİMYA SEKTÖRÜNÜN BİLGİ YÜKÜ
# =======================================================================
s = add_slide()
kicker(s, "Neden Şimdi?")
title(s, "Kimya Sektöründe Bilgi İşi Yükü Büyüyor")
problems = [
    ("⚖", "Regülasyon baskısı",
     "KKDİK, REACH, CLP ve sürekli güncellenen AB/ECHA mevzuatını takip etmek "
     "tam zamanlı bir iş haline geldi; kaçan bir güncelleme ihracatı durdurabiliyor."),
    ("↗", "Hammadde fiyat oynaklığı",
     "Solvent, pigment, polimer ve özel kimyasal fiyatları küresel piyasada hızla "
     "değişiyor; tedarik ve fiyatlama kararları güncel veri istiyor."),
    ("✎", "Dokümantasyon yükü",
     "TDS, spesifikasyon, teklif ve teknik yazışmalar mühendislerin saatlerini alıyor; "
     "Ar-Ge asıl işine zaman ayıramıyor."),
    ("◉", "Pazar ve rakip takibi manuel",
     "Rakip ürün lansmanları, fiyat hareketleri ve pazar trendleri düzensiz ve "
     "kişiye bağlı izleniyor; fırsatlar geç fark ediliyor."),
]
positions = [(Inches(0.7), Inches(2.0)), (Inches(6.75), Inches(2.0)),
             (Inches(0.7), Inches(4.25)), (Inches(6.75), Inches(4.25))]
for (px, py), (ic, hd, bd) in zip(positions, problems):
    card(s, px, py, Inches(5.9), Inches(2.05), ic, hd, bd,
         accent=ACCENT3, head_size=17, body_size=13)

# =======================================================================
# SLIDE 3 — ÇÖZÜM AKIŞI
# =======================================================================
s = add_slide()
kicker(s, "Çözüm", color=ACCENT2)
title(s, "AgentArmy: Görev Verin, Denetlenmiş Çıktı Alın")
txt(s, Inches(0.7), Inches(1.9), Inches(12), Inches(0.8),
    "Soru-cevap aracı değil; hedefi parçalayan, paralel çalışan ve kalite "
    "kapılarından geçen\nkaynaklı çıktı üreten çoklu ajan sistemi.",
    size=17, color=MUTED, line_spacing=1.15)
flow = [("◎", "Hedef", "Doğal dille\nisteğinizi yazın"),
        ("♛", "CEO Ajan", "Görevi parçalar\nve dağıtır"),
        ("⚙", "Ajan Ordusu", "Araştır · Analiz\nYaz · Denetle"),
        ("✓", "Kalite Kapıları", "Kaynak ve tutarlılık\ndenetimi"),
        ("▤", "Çıktı", "Kaynaklı,\ngüvenilir rapor")]
step_flow(s, flow, Inches(3.1))
txt(s, Inches(0.7), Inches(5.7), Inches(12), Inches(1.2),
    "Örnek: “Epoksi reçine pazarında son çeyrek fiyat hareketlerini ve ilk 5 "
    "tedarikçiyi çıkar” →\naraştırma, analiz, rapor ve kaynak doğrulaması "
    "otomatik; siz yalnız sonucu onaylarsınız.",
    size=15, color=WHITE, line_spacing=1.2)

# =======================================================================
# SLIDE 4 — PLATFORM YETENEKLERİ
# =======================================================================
s = add_slide()
kicker(s, "Platform")
title(s, "Kurumsal Kullanım İçin Tasarlanmış Yetenekler")
caps = [
    ("❖", "Playbook orkestrasyonu",
     "Tekrarlanan işler adım adım tanımlanır; her seferinde aynı kalitede, "
     "izlenebilir şekilde çalışır."),
    ("⚙", "Uzman ajan rolleri",
     "Araştırmacı, Analist, Yazar, Editör, Denetçi ve araç kullanan Operatör; "
     "her rolün net sorumluluğu var."),
    ("◎", "Web kaynaklı araştırma",
     "Canlı web aramasıyla güncel veri; her iddia URL kaynağıyla raporun "
     "sonuna eklenir."),
    ("☑", "Kalite kapıları (Verifier)",
     "Denetimsiz final yok: kaynak doğrulama, çelişki yakalama ve rubrik "
     "bazlı kalite kontrolü."),
    ("⚠", "Risk & onay yönetimi",
     "R0–R3 risk sınıfları; sipariş gibi kritik işlemler insan onayı "
     "olmadan asla gerçekleşmez."),
    ("☰", "Kalıcı bilgi tabanı",
     "Her çalışmadan çıkan doğrulanmış bulgular (facts) birikir; "
     "kurumsal hafıza zamanla güçlenir."),
]
xs = [Inches(0.7), Inches(4.72), Inches(8.74)]
ys = [Inches(2.0), Inches(4.15)]
idx = 0
for row in ys:
    for col in xs:
        ic, hd, bd = caps[idx]
        card(s, col, row, Inches(3.85), Inches(1.95), ic, hd, bd,
             accent=ACCENT, head_size=15.5, body_size=12)
        idx += 1
txt(s, Inches(0.7), Inches(6.35), Inches(12), Inches(0.6),
    "Tümü web tabanlı portaldan yönetilir: iş başlatma, canlı izleme, "
    "onay kuyruğu ve raporlar tek ekranda.",
    size=13.5, color=ACCENT, bold=True)

# =======================================================================
# SLIDE 5 — KULLANIM ALANLARI HARİTASI
# =======================================================================
s = add_slide()
kicker(s, "Kullanım Alanları", color=ACCENT2)
title(s, "İşletmenizin Üç Kolu İçin Tek Platform")
areas = [
    ("⚒", "Üretim & Ticaret", ACCENT,
     ["Pazar ve rakip istihbaratı (haftalık brief)",
      "Rakip fiyat karşılaştırma raporları",
      "Ürün lansman ve trend radarı",
      "Teklif ve satış destek içerikleri"]),
    ("⇄", "Tedarik & Stok", ACCENT3,
     ["Stok eşiği izleme (7/24 otomatik)",
      "Gerçek fiyat/tedarikçi araştırması",
      "İnsan onaylı satın alma emri",
      "Kargo takibi ve stok yenileme"]),
    ("⚗", "Ar-Ge & Uyum", ACCENT2,
     ["Regülasyon değişiklik radarı (KKDİK/REACH)",
      "Literatür ve patent tarama özetleri",
      "TDS / teknik doküman taslakları",
      "Müşteri teknik sorularına yanıt taslakları"]),
]
xs = [Inches(0.7), Inches(4.98), Inches(9.26)]
for col, (ic, hd, cc, items) in zip(xs, areas):
    panel(s, col, Inches(2.0), Inches(3.4), Inches(4.3))
    txt(s, col + Inches(0.25), Inches(2.2), Inches(3.0), Inches(0.5),
        ic, size=26, color=cc, bold=True)
    txt(s, col + Inches(0.25), Inches(2.75), Inches(3.0), Inches(0.5),
        hd, size=18, color=WHITE, bold=True)
    bullets(s, col + Inches(0.25), Inches(3.3), Inches(3.0), Inches(2.9),
            items, size=12.5, color=WHITE, gap=8, bullet_color=cc)
txt(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.5),
    "Aynı çekirdek platform; kimya sektörüne özel playbook, persona ve "
    "doğrulama kurallarıyla yapılandırılır — sıfırdan yazılım geliştirme gerekmez.",
    size=13, color=MUTED)

# =======================================================================
# SLIDE 6 — SENARYO 1: PAZAR & RAKİP İSTİHBARATI
# =======================================================================
s = add_slide()
kicker(s, "Senaryo 1 · Üretim & Ticaret")
title(s, "Pazar ve Rakip İstihbaratı — Otomatik Haftalık Brief")
txt(s, Inches(0.7), Inches(1.9), Inches(12), Inches(0.7),
    "Her pazartesi sabahı, ekibiniz daha kahvesini içerken hazır: "
    "kaynaklı, denetlenmiş pazar özeti.",
    size=16, color=MUTED)
left_items = [
    ("Haftalık pazar brifi: ", "hedef segmentte (örn. boya kimyasalları, "
     "yapıştırıcılar) haftanın gelişmeleri, tek sayfa özet."),
    ("Rakip profili: ", "seçilen rakibin ürün gamı, yeni lansmanları, "
     "duyuruları ve stratejik hamleleri."),
    ("Fiyat karşılaştırma: ", "rakip ürün/fiyat verilerinin düzenli "
     "karşılaştırmalı tablosu."),
    ("Bilgi birikimi: ", "her çalışmadan doğrulanmış bulgular kalıcı "
     "bilgi tabanına işlenir; raporlar giderek keskinleşir."),
]
bullets(s, Inches(0.7), Inches(2.7), Inches(6.6), Inches(3.6),
        left_items, size=14.5, gap=12)
panel(s, Inches(7.6), Inches(2.6), Inches(5.0), Inches(3.9))
txt(s, Inches(7.85), Inches(2.85), Inches(4.5), Inches(0.4),
    "NASIL ÇALIŞIR", size=12, color=ACCENT, bold=True)
mini = [
    "1. Zamanlanmış paket (bundle) tetiklenir",
    "2. Araştırmacı canlı web verisi toplar",
    "3. Analist iddiaları ve sayıları test eder",
    "4. Yazar brifi kurar, Editör düzenler",
    "5. Denetçi kaynakları doğrular (PASS/FAIL)",
    "6. Rapor + URL kaynak listesi hazır",
]
bullets(s, Inches(7.85), Inches(3.3), Inches(4.55), Inches(2.9),
        mini, size=13, gap=8, bullet_color=ACCENT2)
txt(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.5),
    "Kazanım: pazarlama/strateji ekibinin haftalar süren manuel taraması, "
    "denetlenmiş otomatik rapora dönüşür.",
    size=13.5, color=WHITE, bold=True)

# =======================================================================
# SLIDE 7 — SENARYO 2: TEDARİK OTOMASYONU
# =======================================================================
s = add_slide()
kicker(s, "Senaryo 2 · Tedarik & Stok", color=ACCENT3)
title(s, "Hammadde Tedarik Otomasyonu — Uçtan Uca")
txt(s, Inches(0.7), Inches(1.85), Inches(12), Inches(0.7),
    "Stok eşiğin altına düşen hammadde, ambalaj veya sarf malzemesi için "
    "sistem kendiliğinden harekete geçer — sipariş ise yalnız sizin onayınızla verilir.",
    size=15, color=MUTED, line_spacing=1.15)
flow2 = [("❶", "Stok İzleme", "Eşik altı tespit\n(15 dk periyot)"),
         ("❷", "Ürün Araştırma", "Gerçek fiyat,\nsatıcı ve link"),
         ("❸", "Karşılaştırma", "En uygun seçenek\n+ link doğrulama"),
         ("❹", "İnsan Onayı", "R3 onay kuyruğu:\nOnayla / Reddet"),
         ("❺", "Sipariş & Kargo", "PO oluşur,\nkargo takibi"),
         ("❻", "Stok Yenileme", "Teslimatta stok\notomatik güncellenir")]
step_flow(s, flow2, Inches(3.0), cw=Inches(1.92), ch=Inches(2.0),
          cgap=Inches(0.14), color=ACCENT3, head_size=13.5, body_size=10.5)
bullets(s, Inches(0.7), Inches(5.5), Inches(12), Inches(1.4), [
    ("Gerçek veri, uydurma yok: ", "fiyat ve linkler gerçek arama "
     "servislerinden gelir; yapay zeka yalnız yorumlar."),
    ("Kontrol sizde: ", "satın alma emri (R3) insan onayı olmadan asla "
     "geçmez; onay ekranında gerekçe ve alternatifler görünür."),
], size=14, gap=8, bullet_color=ACCENT3)

# =======================================================================
# SLIDE 8 — SENARYO 3: REGÜLASYON & UYUM
# =======================================================================
s = add_slide()
kicker(s, "Senaryo 3 · Uyum", color=ACCENT2)
title(s, "Regülasyon Radarı — KKDİK, REACH, CLP Takibi")
cards8 = [
    ("◉", "Değişiklik radarı",
     "ECHA, AB Resmi Gazetesi ve ulusal mevzuat kaynakları düzenli taranır; "
     "sektörünüzü ilgilendiren değişiklikler kaynaklı özet olarak gelir."),
    ("◎", "Etki analizi",
     "Değişikliğin ürün portföyünüze olası etkisi: hangi ürün grubu, "
     "hangi yükümlülük, hangi tarih. Analist + Denetçi zinciriyle çelişkisiz rapor."),
    ("▤", "Uyum doküman desteği",
     "Kayıt dosyaları, bildirim yazışmaları ve iç eğitim özetleri için "
     "taslak hazırlığı; nihai karar her zaman uzmanınızda."),
    ("◷", "Zamanlanmış çalışma",
     "Haftalık/aylık otomatik tarama; yeni bir şey yoksa kısa 'değişiklik yok' "
     "notu, varsa detaylı brif. Ekibiniz yalnız aksiyon gerektirenle ilgilenir."),
]
positions = [(Inches(0.7), Inches(2.0)), (Inches(6.75), Inches(2.0)),
             (Inches(0.7), Inches(4.2)), (Inches(6.75), Inches(4.2))]
for (px, py), (ic, hd, bd) in zip(positions, cards8):
    card(s, px, py, Inches(5.9), Inches(2.0), ic, hd, bd,
         accent=ACCENT2, head_size=16, body_size=12.5)
txt(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.6),
    "Not: Platform hukuki danışmanlık vermez; uzmanlarınızın kararını "
    "hızlandıran, kaynaklı ve denetlenmiş ön çalışma üretir.",
    size=12.5, color=MUTED, italic=True)

# =======================================================================
# SLIDE 9 — SENARYO 4: AR-GE & TEKNİK DOKÜMANTASYON
# =======================================================================
s = add_slide()
kicker(s, "Senaryo 4 · Ar-Ge & Laboratuvar")
title(s, "Ar-Ge Ekibinize Dijital Asistan Ordusu")
left9 = [
    ("Literatür & patent tarama: ", "yeni formülasyon çalışmaları öncesi "
     "mevcut yaklaşımların kaynaklı özeti; günler süren tarama saatlere iner."),
    ("Alternatif hammadde araştırması: ", "kısıtlanan veya pahalılaşan bir "
     "girdiye ikame seçenekleri, tedarikçi ve fiyat bilgisiyle."),
    ("TDS / teknik doküman taslakları: ", "ürün verilerinizden standart "
     "formatta taslak; Editör dil ve format tutarlılığını sağlar."),
    ("Müşteri teknik soruları: ", "sık gelen uygulama/uyumluluk sorularına "
     "kaynaklı yanıt taslakları; mühendis yalnız kontrol edip gönderir."),
    ("Rapor standartlaştırma: ", "deney ve proje raporlarının tek formatta, "
     "eksiksiz bölümlerle derlenmesi."),
]
bullets(s, Inches(0.7), Inches(2.1), Inches(7.2), Inches(4.2),
        left9, size=14, gap=12)
panel(s, Inches(8.2), Inches(2.1), Inches(4.4), Inches(4.2))
txt(s, Inches(8.45), Inches(2.35), Inches(3.9), Inches(0.4),
    "NEDEN GÜVENİLİR?", size=12, color=ACCENT, bold=True)
bullets(s, Inches(8.45), Inches(2.8), Inches(3.95), Inches(3.3), [
    "Her iddia URL kaynağıyla gelir",
    "Denetçi ajan çelişkileri yakalar",
    "Kimyaya özel doğrulama rubriği",
    "Nihai onay her zaman insanda",
    "Tüm adımlar kayıt altında (denetim izi)",
], size=13, gap=10, bullet_color=ACCENT2)
txt(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.5),
    "Kazanım: Ar-Ge mühendisleri doküman işçiliğinden kurtulur, "
    "laboratuvara ve inovasyona odaklanır.",
    size=13.5, color=WHITE, bold=True)

# =======================================================================
# SLIDE 10 — PORTAL
# =======================================================================
s = add_slide()
kicker(s, "Yönetim", color=ACCENT3)
title(s, "Web Portal: Her Şey Tek Ekranda, Kod Gerekmez")
caps10 = [
    ("▣", "İş başlatma & canlı izleme",
     "Doğal dille yeni iş açın; adımları ve çıktıları canlı izleyin."),
    ("▦", "Stok ekranı",
     "Ürün ekleyin, eşik ve hedef belirleyin; izlemeyi açıp kapatın."),
    ("☑", "Onay kuyruğu",
     "Bekleyen kritik işlemler için satır içi Onayla/Reddet + gerekçe paneli."),
    ("▤", "Tedarik raporu",
     "Stok tetikleri, siparişler ve kargo durumu; '● Canlı' otomatik yenileme."),
    ("❖", "İçerik yönetimi",
     "Playbook, persona, paket ve araç tanımları ekranlardan yönetilir."),
    ("∑", "KPI & raporlama",
     "Tamamlanan operasyonlardan otomatik KPI özetleri ve dışa aktarım."),
]
xs = [Inches(0.7), Inches(4.72), Inches(8.74)]
ys = [Inches(2.0), Inches(3.95)]
idx = 0
for row in ys:
    for col in xs:
        ic, hd, bd = caps10[idx]
        card(s, col, row, Inches(3.85), Inches(1.75), ic, hd, bd,
             accent=ACCENT3, head_size=15, body_size=12)
        idx += 1
txt(s, Inches(0.7), Inches(6.1), Inches(12), Inches(0.8),
    "Satın alma, kalite, Ar-Ge ve yönetim aynı platformda çalışır; "
    "yetkiler role göre ayrışır.\nKurulum sonrası yeni süreç eklemek "
    "yazılım projesi değil, yapılandırma işidir.",
    size=13.5, color=MUTED, line_spacing=1.2)

# =======================================================================
# SLIDE 11 — GÜVENLİK & KONTROL
# =======================================================================
s = add_slide()
kicker(s, "Güven & Kontrol")
title(s, "Özerklik Artar, Kontrol Sizde Kalır")
risk_rows = [
    ("R0  ·  Okuma", "Veri okuma, araştırma, raporlama — serbest çalışır.", ACCENT2),
    ("R1  ·  Düşük riskli yazma", "Stok güncelleme gibi geri alınabilir işlemler — kayıtlı ve izlenebilir.", ACCENT),
    ("R2  ·  Hassas işlem", "Dış sistemlere veri gönderimi — onay kuyruğuna düşer.", ACCENT3),
    ("R3  ·  Kritik işlem", "Satın alma emri gibi mali sonuçlu işlemler — insan onayı zorunlu, onaysız asla geçmez.", RGBColor(0xC0, 0x39, 0x2B)),
]
ly = Inches(2.0)
for hd, bd, cc in risk_rows:
    panel(s, Inches(0.7), ly, Inches(7.4), Inches(0.95))
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), ly,
                             Inches(0.12), Inches(0.95))
    bar.fill.solid(); bar.fill.fore_color.rgb = cc
    bar.line.fill.background(); bar.shadow.inherit = False
    txt(s, Inches(1.05), ly + Inches(0.12), Inches(6.9), Inches(0.4), hd,
        size=15, color=WHITE, bold=True)
    txt(s, Inches(1.05), ly + Inches(0.5), Inches(6.9), Inches(0.4), bd,
        size=12, color=MUTED)
    ly = Emu(int(ly) + int(Inches(1.1)))
panel(s, Inches(8.5), Inches(2.0), Inches(4.1), Inches(4.3))
txt(s, Inches(8.75), Inches(2.25), Inches(3.6), Inches(0.4),
    "AYRICA", size=12, color=ACCENT, bold=True)
bullets(s, Inches(8.75), Inches(2.7), Inches(3.65), Inches(3.4), [
    "Tüm adımların denetim izi (kim, ne, ne zaman)",
    "Çıktılarda zorunlu kaynak listesi",
    "Rol bazlı erişim ve izin matrisi",
    "Maliyet takibi: model kullanımı iş bazında izlenir",
    "Verileriniz size ait; satıcı kilidi yok",
], size=12.5, gap=10)
txt(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.5),
    "Felsefemiz: “denetimsiz final yok” — hızı otomasyondan, "
    "güveni insan onayından alırız.",
    size=13.5, color=WHITE, bold=True)

# =======================================================================
# SLIDE 12 — KAZANIMLAR / ROI
# =======================================================================
s = add_slide()
kicker(s, "İş Değeri", color=ACCENT2)
title(s, "Beklenen Kazanımlar")
stats = [
    ("%60–80", "zaman tasarrufu", "pazar araştırması, brif ve rapor "
     "hazırlığında manuel eforun otomasyona devri"),
    ("7/24", "kesintisiz izleme", "stok, fiyat ve regülasyon takibi mesai "
     "saatine bağlı olmaktan çıkar"),
    ("%100", "kaynaklı çıktı", "her rapor URL kaynakları ve denetçi "
     "kontrolüyle gelir; “kaynağı neydi?” sorusu biter"),
]
xs = [Inches(0.7), Inches(4.98), Inches(9.26)]
for col, (big, lab, desc) in zip(xs, stats):
    panel(s, col, Inches(2.0), Inches(3.4), Inches(2.6))
    txt(s, col, Inches(2.3), Inches(3.4), Inches(0.9), big, size=44,
        color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    txt(s, col, Inches(3.25), Inches(3.4), Inches(0.4), lab, size=15,
        color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, col + Inches(0.3), Inches(3.7), Inches(2.8), Inches(0.8), desc,
        size=11.5, color=MUTED, align=PP_ALIGN.CENTER, line_spacing=1.1)
txt(s, Inches(0.7), Inches(5.0), Inches(12), Inches(0.4),
    "Örnek hesap (varsayımsal):", size=14, color=WHITE, bold=True)
bullets(s, Inches(0.7), Inches(5.45), Inches(12), Inches(1.3), [
    ("Pazar/rakip takibi: ", "haftada 2 kişi × 1 gün → otomatik brif + "
     "30 dk kontrol ≈ ayda ~14 iş günü kazanım."),
    ("Tedarik araştırması: ", "eşik altı her ürün için ~yarım günlük "
     "araştırma → dakikalar içinde onaya hazır öneri."),
], size=13.5, gap=8, bullet_color=ACCENT2)
txt(s, Inches(0.7), Inches(6.8), Inches(12), Inches(0.4),
    "* Oranlar örnek senaryolara dayalı hedeflerdir; pilot çalışmada "
    "kendi süreçlerinizle birlikte ölçülür.",
    size=11, color=MUTED, italic=True)

# =======================================================================
# SLIDE 13 — PİLOT PLANI
# =======================================================================
s = add_slide()
kicker(s, "Yol Haritası")
title(s, "4 Adımda Pilot: Düşük Risk, Hızlı Sonuç")
phases = [
    ("1", "Keşif", "1–2 hafta",
     "Süreç ve ağrı haritası, terminoloji sözlüğü, hedef senaryoların "
     "seçimi (örn. tedarik + regülasyon radarı)."),
    ("2", "Yapılandırma", "1–2 hafta",
     "Kimyaya özel playbook, persona ve doğrulama kuralları; stok/araç "
     "entegrasyonları; onay akışlarının kurulumu."),
    ("3", "Pilot çalışma", "3–4 hafta",
     "Seçilen ekip gerçek işlerde kullanır; haftalık geri bildirim "
     "döngüsü; kalite ve süre metrikleri toplanır."),
    ("4", "Değerlendirme", "1 hafta",
     "KPI raporu: zaman kazanımı, kalite, onay istatistikleri. "
     "Yaygınlaştırma kararı ve yol haritası."),
]
xs = [Inches(0.7), Inches(3.88), Inches(7.06), Inches(10.24)]
for col, (no, hd, dur, bd) in zip(xs, phases):
    panel(s, col, Inches(2.1), Inches(2.95), Inches(3.6))
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, col + Inches(0.25),
                              Inches(2.35), Inches(0.55), Inches(0.55))
    circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT
    circ.line.fill.background(); circ.shadow.inherit = False
    txt(s, col + Inches(0.25), Inches(2.44), Inches(0.55), Inches(0.4), no,
        size=20, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True,
        align=PP_ALIGN.CENTER)
    txt(s, col + Inches(0.95), Inches(2.38), Inches(1.9), Inches(0.4), hd,
        size=17, color=WHITE, bold=True)
    txt(s, col + Inches(0.95), Inches(2.72), Inches(1.9), Inches(0.35), dur,
        size=12, color=ACCENT3, bold=True)
    txt(s, col + Inches(0.25), Inches(3.25), Inches(2.45), Inches(2.3), bd,
        size=12, color=MUTED, line_spacing=1.15)
txt(s, Inches(0.7), Inches(6.1), Inches(12), Inches(0.8),
    "Toplam ~6–9 hafta içinde ölçülmüş sonuçlarla karar noktasına "
    "gelirsiniz.\nPilot kapsamı dar tutulur: 1–2 senaryo, seçili ekip, net başarı kriterleri.",
    size=14, color=WHITE, line_spacing=1.2)

# =======================================================================
# SLIDE 14 — KAPANIŞ
# =======================================================================
s = add_slide()
strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.25), SLIDE_H)
strip.fill.solid(); strip.fill.fore_color.rgb = ACCENT
strip.line.fill.background(); strip.shadow.inherit = False
txt(s, Inches(0.9), Inches(1.1), Inches(11), Inches(0.5),
    "SONRAKİ ADIMLAR", size=15, color=ACCENT, bold=True)
txt(s, Inches(0.9), Inches(1.8), Inches(11.5), Inches(1.6),
    "Kimya sektöründeki bilgi işinizi\nbirlikte otomatikleştirelim.",
    size=40, color=WHITE, bold=True, line_spacing=1.05)
steps = [
    ("1. Keşif toplantısı (60 dk): ", "öncelikli süreçlerinizi ve ağrı "
     "noktalarını birlikte haritalayalım."),
    ("2. Canlı demo: ", "tedarik otomasyonu ve pazar brifi akışlarını "
     "gerçek portal üzerinde görün."),
    ("3. Pilot teklifi: ", "kapsam, başarı kriterleri ve takvimi içeren "
     "pilot planını 1 hafta içinde sunalım."),
]
bullets(s, Inches(0.9), Inches(3.7), Inches(11), Inches(2.0),
        steps, size=17, gap=14)
panel(s, Inches(0.9), Inches(5.9), Inches(11.5), Inches(0.85))
txt(s, Inches(1.2), Inches(6.12), Inches(11), Inches(0.5),
    "AgentArmy  ·  Kurumsal Yapay Zeka Ajan Platformu  ·  "
    "kadirogreten@hotmail.com",
    size=14, color=WHITE, bold=True)

prs.save("AgentArmy-Kimya-Sunumu.pptx")
print("OK: AgentArmy-Kimya-Sunumu.pptx")
