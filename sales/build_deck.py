#!/usr/bin/env python3
"""
AgentArmy satış sunumu (PPTX) üretici.
Kurumsal, koyu temalı, satış odaklı bir deste oluşturur.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- Marka renk paleti -------------------------------------------------
# --- AYDINLIK TEMA (teal + sıcak turuncu) ---
# Not: WHITE artık "ana koyu metin" anlamında kullanılıyor (aydınlık zeminde).
BG_DARK   = RGBColor(0xF3, 0xF6, 0xFA)   # ferah açık zemin
BG_PANEL  = RGBColor(0xFF, 0xFF, 0xFF)   # beyaz kartlar
ACCENT    = RGBColor(0x0E, 0x9C, 0x96)   # teal (birincil vurgu)
ACCENT2   = RGBColor(0x2F, 0xA8, 0x55)   # yeşil
ACCENT3   = RGBColor(0xE2, 0x73, 0x1F)   # sıcak turuncu
WHITE     = RGBColor(0x16, 0x22, 0x32)   # ana koyu metin (lacivert-antrasit)
MUTED     = RGBColor(0x57, 0x64, 0x78)   # ikincil gri metin
CARD_LINE = RGBColor(0xDD, 0xE4, 0xEF)   # açık kart kenarlığı
CHIP_INK  = RGBColor(0xFF, 0xFF, 0xFF)   # parlak çip üstü beyaz metin

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT = "Calibri"
FONT_H = "Calibri"


def fix(text):
    r"""Emoji surrogate ciftlerini (orn. \uD83C\uDFAF) tek karaktere birlestirir;
    eslesmeyen kalinti surrogate'leri temizler."""
    try:
        text = text.encode("utf-16", "surrogatepass").decode("utf-16")
    except Exception:
        pass
    return "".join(c for c in text if not (0xD800 <= ord(c) <= 0xDFFF))

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_DARK
    bg.line.fill.background()
    bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element)
    s.shapes._spTree.insert(2, bg._element)
    return s


def txt(slide, left, top, width, height, text, size=18, color=WHITE,
        bold=False, align=PP_ALIGN.LEFT, font=FONT, anchor=MSO_ANCHOR.TOP,
        line_spacing=1.0, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = fix(text).split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
        r.font.color.rgb = color
    return tb


def bullets(slide, left, top, width, height, items, size=18, color=WHITE,
            gap=6, bullet_color=ACCENT, bold_lead=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        rb = p.add_run()
        rb.text = "\u25B8  "
        rb.font.size = Pt(size)
        rb.font.color.rgb = bullet_color
        rb.font.name = FONT
        rb.font.bold = True
        if isinstance(it, tuple):
            lead, rest = it
            r1 = p.add_run(); r1.text = fix(lead)
            r1.font.size = Pt(size); r1.font.bold = True
            r1.font.color.rgb = color; r1.font.name = FONT
            r2 = p.add_run(); r2.text = fix(rest)
            r2.font.size = Pt(size); r2.font.color.rgb = MUTED
            r2.font.name = FONT
        else:
            r = p.add_run(); r.text = fix(it)
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.name = FONT; r.font.bold = bold_lead
    return tb


def panel(slide, left, top, width, height, fill=BG_PANEL, line=CARD_LINE,
          radius=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def accent_bar(slide, left, top, width=Inches(0.9), color=ACCENT):
    b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Pt(6))
    b.fill.solid(); b.fill.fore_color.rgb = color
    b.line.fill.background(); b.shadow.inherit = False
    return b


def kicker(slide, text, color=ACCENT):
    txt(slide, Inches(0.7), Inches(0.45), Inches(8), Inches(0.4),
        text.upper(), size=13, color=color, bold=True)


def title(slide, text, size=30):
    txt(slide, Inches(0.7), Inches(0.8), Inches(12), Inches(1.0),
        text, size=size, color=WHITE, bold=True)
    accent_bar(slide, Inches(0.72), Inches(1.62))


def card(slide, left, top, width, height, icon, head, body,
         accent=ACCENT, head_size=16, body_size=12.5):
    p = panel(slide, left, top, width, height)
    txt(slide, left + Inches(0.25), top + Inches(0.18), width - Inches(0.5),
        Inches(0.5), icon, size=22, color=accent, bold=True)
    txt(slide, left + Inches(0.25), top + Inches(0.62), width - Inches(0.5),
        Inches(0.5), head, size=head_size, color=WHITE, bold=True)
    txt(slide, left + Inches(0.25), top + Inches(1.05), width - Inches(0.5),
        height - Inches(1.15), body, size=body_size, color=MUTED,
        line_spacing=1.08)
    return p


# =======================================================================
# SLIDE 1 — KAPAK
# =======================================================================
s = add_slide()
# dekoratif paneller
strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.25), SLIDE_H)
strip.fill.solid(); strip.fill.fore_color.rgb = ACCENT
strip.line.fill.background(); strip.shadow.inherit = False

txt(s, Inches(0.9), Inches(1.1), Inches(11), Inches(0.5),
    "AGENTARMY  \u00B7  KURUMSAL YAPAY ZEKA AJAN PLATFORMU", size=15,
    color=ACCENT, bold=True)
txt(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(2.0),
    "Sizin Yerinize \u00C7al\u0131\u015fan\nYapay Zeka Ajan Ordusu", size=46,
    color=WHITE, bold=True, line_spacing=1.0)
txt(s, Inches(0.9), Inches(4.0), Inches(10.5), Inches(1.2),
    "Soru-cevap de\u011fil, g\u00f6rev \u2192 \u00e7\u0131kt\u0131. Kan\u0131t, izlenebilirlik ve kontrol\n"
    "varsay\u0131lan olarak gelir. \u00c7oklu ajan ekipleriyle paralel \u00fcretim.",
    size=18, color=MUTED, line_spacing=1.2)

# alt rozetler
badges = ["\u00C7oklu Ajan Orkestrasyonu", "Kalite Kap\u0131lar\u0131", "Risk & Onay Y\u00f6netimi", "Web Tabanl\u0131 Portal"]
bx = Inches(0.9)
for b in badges:
    w = Inches(2.95)
    pn = panel(s, bx, Inches(5.55), w, Inches(0.6), fill=BG_PANEL)
    txt(s, bx, Inches(5.68), w, Inches(0.4), b, size=12.5, color=WHITE,
        bold=True, align=PP_ALIGN.CENTER)
    bx = Emu(int(bx) + int(w) + int(Inches(0.12)))

txt(s, Inches(0.9), Inches(6.7), Inches(11), Inches(0.4),
    "M\u00fc\u015fteri Sunumu  \u00B7  Gizli ve \u00d6zeldir", size=12, color=MUTED)

# =======================================================================
# SLIDE 2 — SORUN
# =======================================================================
s = add_slide()
kicker(s, "Neden \u015eimdi?")
title(s, "\u015eirketler Bilgi \u0130\u015finde Bo\u011fuluyor")
problems = [
    ("\u23F1", "Zamana yenik d\u00fc\u015f\u00fcyor",
     "Pazar ara\u015ft\u0131rmas\u0131, rapor ve analiz manuel; g\u00fcnler s\u00fcr\u00fcyor, ekip kapasitesi t\u0131kan\u0131yor."),
    ("\u2753", "Kalite tutars\u0131z",
     "Her \u00e7\u0131kt\u0131 ki\u015fiye ba\u011fl\u0131; kaynaks\u0131z iddialar, \u00e7eli\u015fkiler ve format kaymalar\u0131 olu\u015fuyor."),
    ("\uD83D\uDD0D", "\u0130zlenemiyor",
     "\u00c7\u0131kt\u0131 nas\u0131l \u00fcretildi, hangi kaynaktan geldi belli de\u011fil; denetim ve g\u00fcven zor."),
    ("\uD83D\uDCB8", "Genel AI ara\u00e7lar\u0131 riskli",
     "Kontrolsuz otomasyon, maliyet patlamas\u0131 ve y\u00f6neti\u015fim eksikli\u011fi kurumsal kullan\u0131m\u0131 zorla\u015ft\u0131r\u0131yor."),
]
gx = Inches(0.7); gw = Inches(5.9); gh = Inches(2.05)
positions = [(Inches(0.7), Inches(2.0)), (Inches(6.75), Inches(2.0)),
             (Inches(0.7), Inches(4.25)), (Inches(6.75), Inches(4.25))]
colors = [ACCENT3, ACCENT3, ACCENT3, ACCENT3]
for (px, py), (ic, hd, bd), cc in zip(positions, problems, colors):
    card(s, px, py, gw, gh, ic, hd, bd, accent=cc, head_size=18, body_size=14)

# =======================================================================
# SLIDE 3 — ÇÖZÜM
# =======================================================================
s = add_slide()
kicker(s, "\u00c7\u00f6z\u00fcm", color=ACCENT2)
title(s, "AgentArmy: \u00c7oklu Ajan \u00dcretim Sistemi")
txt(s, Inches(0.7), Inches(1.9), Inches(12), Inches(0.8),
    "Bir hedef verirsiniz \u2014 CEO ajan i\u015fi par\u00e7alar, uzman ajanlar paralel \u00e7al\u0131\u015f\u0131r,\n"
    "kalite kap\u0131lar\u0131ndan ge\u00e7en kaynakl\u0131 \u00e7\u0131kt\u0131y\u0131 size sunar.",
    size=17, color=MUTED, line_spacing=1.15)

flow = [("\uD83C\uDFAF", "Hedef", "Do\u011fal dille\nisте\u011finizi yaz\u0131n"),
        ("\uD83D\uDC54", "CEO Ajan", "G\u00f6revi par\u00e7alar\nve da\u011f\u0131t\u0131r"),
        ("\uD83E\uDD16", "Ajan Ordusu", "Ara\u015ft\u0131r \u00B7 Analiz\nYaz \u00B7 Denetle"),
        ("\u2705", "Kalite Kap\u0131lar\u0131", "QG1\u2013QG5\ndenetimi"),
        ("\uD83D\uDCC4", "\u00c7\u0131kt\u0131", "Kaynakl\u0131,\ng\u00fcvenilir rapor")]
n = len(flow)
cw = Inches(2.18); cgap = Inches(0.18)
totw = n * int(cw) + (n - 1) * int(cgap)
startx = int((int(SLIDE_W) - totw) / 2)
yy = Inches(3.1)
for i, (ic, hd, bd) in enumerate(flow):
    cx = Emu(startx + i * (int(cw) + int(cgap)))
    pn = panel(s, cx, yy, cw, Inches(2.1), fill=BG_PANEL)
    txt(s, cx, yy + Inches(0.25), cw, Inches(0.6), ic, size=30,
        color=ACCENT2, align=PP_ALIGN.CENTER)
    txt(s, cx, yy + Inches(0.95), cw, Inches(0.5), hd, size=16,
        color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, cx, yy + Inches(1.4), cw, Inches(0.6), bd, size=12,
        color=MUTED, align=PP_ALIGN.CENTER, line_spacing=1.05)
    if i < n - 1:
        ar = s.shapes.add_shape(MSO_SHAPE.CHEVRON,
            Emu(cx + int(cw) - int(Inches(0.02))), yy + Inches(0.78),
            Inches(0.22), Inches(0.5))
        ar.fill.solid(); ar.fill.fore_color.rgb = ACCENT
        ar.line.fill.background(); ar.shadow.inherit = False

txt(s, Inches(0.7), Inches(5.7), Inches(12), Inches(1.2),
    "Sonu\u00e7: \u201cH\u0131zl\u0131 ama da\u011f\u0131n\u0131k\u201d de\u011fil \u2014 s\u00fcre\u00e7 + kalite ile \u00fcretim sistemi.\n"
    "\u00d6zerklik artar, riskler izin matrisi ve onay kap\u0131lar\u0131yla kontrol alt\u0131nda kal\u0131r.",
    size=15, color=WHITE, line_spacing=1.2)

# =======================================================================
# SLIDE 4 — MİMARİ / NASIL ÇALIŞIR
# =======================================================================
s = add_slide()
kicker(s, "Mimari")
title(s, "D\u00f6rt Katmanl\u0131, \u00d6l\u00e7eklenebilir Tasar\u0131m")
layers = [
    ("01  Y\u00f6neti\u015fim Katman\u0131", "\u0130zinler, onay ak\u0131\u015flar\u0131, denetim logu, maliyet defteri, geri alma.", ACCENT3),
    ("02  Orkestrasyon Katman\u0131", "CEO ajan g\u00f6revi par\u00e7alar, paralel y\u00fcr\u00fct\u00fcr, \u00e7at\u0131\u015fmalar\u0131 \u00e7\u00f6zer.", ACCENT),
    ("03  Ajan Katman\u0131", "Rol + ara\u00e7lar + bellek + \u00e7al\u0131\u015fma d\u00f6ng\u00fcs\u00fc (planla\u2192\u00fcret\u2192kontrol\u2192d\u00fczelt).", ACCENT2),
    ("04  Zeka Katman\u0131 (LLM)", "OpenAI GPT-4.1 / GPT-5, web_search ile kaynakl\u0131 grounding.", RGBColor(0xB0,0x8CFF&0xFF,0xFF)),
]
ly = Inches(2.0)
for i, (hd, bd, cc) in enumerate(layers):
    pn = panel(s, Inches(0.7), ly, Inches(11.9), Inches(1.05))
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), ly,
                             Inches(0.12), Inches(1.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = cc
    bar.line.fill.background(); bar.shadow.inherit = False
    txt(s, Inches(1.05), ly + Inches(0.16), Inches(4.5), Inches(0.5), hd,
        size=18, color=WHITE, bold=True)
    txt(s, Inches(1.05), ly + Inches(0.6), Inches(11), Inches(0.4), bd,
        size=13.5, color=MUTED)
    ly = Emu(int(ly) + int(Inches(1.22)))

txt(s, Inches(0.7), Inches(6.95), Inches(12), Inches(0.4),
    "Tek Ger\u00e7ek Kayna\u011f\u0131 (Single Source of Truth): Facts \u00B7 Decisions \u00B7 Work depolar\u0131 t\u00fcm ajanlarca payla\u015f\u0131l\u0131r.",
    size=12.5, color=ACCENT, bold=True)

# =======================================================================
# SLIDE 5 — AJAN ROLLERİ
# =======================================================================
s = add_slide()
kicker(s, "Ajan Ordusu", color=ACCENT2)
title(s, "Uzman Roller, Net Sorumluluklar")
roles = [
    ("\uD83D\uDD0E", "Ara\u015ft\u0131rmac\u0131", "Kaynak tarar, al\u0131nt\u0131 ve link \u00e7\u0131kar\u0131r, g\u00fcven puan\u0131 verir."),
    ("\uD83D\uDCCA", "Analist", "\u0130ddialar\u0131 test eder, tutarl\u0131l\u0131k ve say\u0131 kontrol\u00fc yapar."),
    ("\u270D", "Yazar", "Nihai raporu kurar; yap\u0131, ak\u0131\u015f ve arg\u00fcman \u00fcretir."),
    ("\uD83D\uDCDD", "Edit\u00f6r", "Dil, ton, okunabilirlik ve format standard\u0131."),
    ("\uD83D\uDEE1", "Denet\u00e7i", "Kaynak do\u011frular, \u00e7eli\u015fki yakalar, risk etiketler."),
    ("\u2699", "Operat\u00f6r", "Ara\u00e7 \u00e7a\u011f\u0131r\u0131r: sat\u0131n alma, stok, kargo, dosya, web."),
]
cw = Inches(3.85); ch = Inches(1.75)
xs = [Inches(0.7), Inches(4.72), Inches(8.74)]
ys = [Inches(2.0), Inches(3.9)]
idx = 0
for row in ys:
    for col in xs:
        ic, hd, bd = roles[idx]
        card(s, col, row, cw, ch, ic, hd, bd, accent=ACCENT2,
             head_size=17, body_size=13)
        idx += 1
txt(s, Inches(0.7), Inches(5.95), Inches(12), Inches(1.0),
    "Opsiyonel uzmanlar: Contrarian (kar\u015f\u0131 g\u00f6r\u00fc\u015f), Cost/Time Planner, Policy/Privacy.\n"
    "Persona ajanlar\u0131 (PM, Finans, Pazar Ara\u015ft\u0131rmac\u0131s\u0131...) bu \u00e7ekirde\u011fi koordine eder.",
    size=14, color=MUTED, line_spacing=1.15)

# =======================================================================
# SLIDE 6 — KURUMSAL ÖZELLİKLER (PORTAL)
# =======================================================================
s = add_slide()
kicker(s, "\u00dcr\u00fcn")
title(s, "Web Portal \u2014 Tek Ekrandan Y\u00f6netim")
feat = [
    ("\uD83D\uDCC8", "CEO Dashboard", "Ba\u015far\u0131 oran\u0131, \u00e7al\u0131\u015fma ve i\u015f metrikleri, rapor indirme."),
    ("\uD83E\uDDD9", "CEO Sihirbaz\u0131", "Hedef yaz, otomatik plan + soru-cevap, onayla \u00e7al\u0131\u015ft\u0131r."),
    ("\uD83D\uDCB0", "Maliyet Defteri", "Token/\u00fccret takibi, i\u015f ba\u015f\u0131na maliyet g\u00f6r\u00fcn\u00fcrl\u00fc\u011f\u00fc."),
    ("\u2705", "Onay Kuyru\u011fu", "Y\u00fcksek riskli ad\u0131mlar i\u00e7in insan onay\u0131 + geri alma."),
    ("\uD83D\uDCDC", "Denetim Logu", "Her eylemin izlenebilir kayd\u0131, tam \u015feffafl\u0131k."),
    ("\uD83D\uDCDA", "Bilgi Taban\u0131 (Facts)", "Kaynakl\u0131, yeniden kullan\u0131labilir kal\u0131c\u0131 haf\u0131za."),
    ("\uD83E\uDDE9", "Domain Pack & Playbook", "Sekt\u00f6re \u00f6zel ajan, ara\u00e7 ve i\u015f ak\u0131\u015f\u0131 paketleri."),
    ("\uD83D\uDD52", "Zamanlanm\u0131\u015f \u0130\u015fler", "Cron tabanl\u0131 otomatik \u00e7al\u0131\u015ft\u0131rma (\u00f6r. haftal\u0131k brief)."),
]
cw = Inches(2.92); ch = Inches(1.95)
xs = [Inches(0.7), Inches(3.78), Inches(6.86), Inches(9.94)]
ys = [Inches(2.0), Inches(4.1)]
idx = 0
accs = [ACCENT, ACCENT2, ACCENT3]
for row in ys:
    for col in xs:
        ic, hd, bd = feat[idx]
        card(s, col, row, cw, ch, ic, hd, bd, accent=accs[idx % 3],
             head_size=14.5, body_size=12)
        idx += 1

# =======================================================================
# SLIDE 7 — GÜVENLİK & YÖNETİŞİM
# =======================================================================
s = add_slide()
kicker(s, "Kurumsal G\u00fcven", color=ACCENT3)
title(s, "Kontroll\u00fc \u00d6zerklik: Risk & Kalite")
# Sol: Risk seviyeleri
panel(s, Inches(0.7), Inches(2.0), Inches(5.85), Inches(4.6))
txt(s, Inches(1.0), Inches(2.2), Inches(5.3), Inches(0.5),
    "Risk & \u0130zin Matrisi", size=18, color=WHITE, bold=True)
risks = [
    ("R0", "Zarars\u0131z \u2014 tam otomatik, sadece log", ACCENT2),
    ("R1", "D\u00fc\u015f\u00fck \u2014 otomatik + QG1\u2013QG4", ACCENT),
    ("R2", "Orta \u2014 denet\u00e7i onay\u0131 + gerek\u00e7e", ACCENT3),
    ("R3", "Y\u00fcksek \u2014 insan onay\u0131 + geri alma plan\u0131", RGBColor(0xFF,0x6B,0x6B)),
]
ry = Inches(2.85)
for code, desc, cc in risks:
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), ry,
                              Inches(0.85), Inches(0.55))
    chip.fill.solid(); chip.fill.fore_color.rgb = cc
    chip.line.fill.background(); chip.shadow.inherit = False
    txt(s, Inches(1.0), ry + Inches(0.1), Inches(0.85), Inches(0.4), code,
        size=16, color=CHIP_INK, bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(2.0), ry + Inches(0.12), Inches(4.4), Inches(0.5), desc,
        size=13.5, color=WHITE)
    ry = Emu(int(ry) + int(Inches(0.92)))

# Sağ: Kalite kapıları
panel(s, Inches(6.75), Inches(2.0), Inches(5.85), Inches(4.6))
txt(s, Inches(7.05), Inches(2.2), Inches(5.3), Inches(0.5),
    "Kalite Kap\u0131lar\u0131 (Quality Gates)", size=18, color=WHITE, bold=True)
gates = [
    ("QG1", "Yap\u0131 / Kapsam \u2014 do\u011fru format, kapsam d\u0131\u015f\u0131na \u00e7\u0131kmama"),
    ("QG2", "Kan\u0131t / Kaynak \u2014 kritik iddialar\u0131n kayna\u011f\u0131"),
    ("QG3", "Tutarl\u0131l\u0131k \u2014 i\u00e7 \u00e7eli\u015fki ve say\u0131 kontrol\u00fc"),
    ("QG4", "Risk Etiketleme \u2014 belirsizliklerin i\u015faretlenmesi"),
    ("QG5", "\u0130leti\u015fim \u2014 ton, a\u00e7\u0131kl\u0131k, hedef kitle uygunlu\u011fu"),
]
gy = Inches(2.85)
for code, desc in gates:
    txt(s, Inches(7.05), gy, Inches(1.0), Inches(0.4), code, size=14,
        color=ACCENT2, bold=True)
    txt(s, Inches(8.0), gy, Inches(4.4), Inches(0.6), desc, size=13,
        color=MUTED, line_spacing=1.0)
    gy = Emu(int(gy) + int(Inches(0.72)))

# =======================================================================
# SLIDE 8 — KULLANIM ALANLARI
# =======================================================================
s = add_slide()
kicker(s, "Kullan\u0131m Alanlar\u0131", color=ACCENT2)
title(s, "Hangi \u0130\u015fleri Devral\u0131r?")
uses = [
    ("\uD83D\uDCC8", "Pazar Zekas\u0131", "Rakip analizi, trend radar\u0131, fiyat kar\u015f\u0131la\u015ft\u0131rma, haftal\u0131k brief, pazar boyutland\u0131rma."),
    ("\uD83D\uDCCB", "\u00dcr\u00fcn & Strateji", "PRD tasla\u011f\u0131, build-vs-buy, karar notu, yol haritas\u0131 girdileri."),
    ("\uD83D\uDD27", "Teknik", "Teknik tasar\u0131m dok\u00fcman\u0131, PoC plan\u0131, mimari trade-off analizi."),
    ("\uD83C\uDFED", "Operasyon & Tedarik", "Stok kontrol\u00fc, sat\u0131n alma sipari\u015fi, kargo takibi (ara\u00e7 entegrasyonlar\u0131)."),
]
cw = Inches(5.9); ch = Inches(2.05)
positions = [(Inches(0.7), Inches(2.0)), (Inches(6.75), Inches(2.0)),
             (Inches(0.7), Inches(4.25)), (Inches(6.75), Inches(4.25))]
for (px, py), (ic, hd, bd) in zip(positions, uses):
    card(s, px, py, cw, ch, ic, hd, bd, accent=ACCENT, head_size=18, body_size=14)

# =======================================================================
# SLIDE 9 — RAKİPLERDEN FARK
# =======================================================================
s = add_slide()
kicker(s, "Fark\u0131m\u0131z")
title(s, "Genel AI Ara\u00e7lar\u0131ndan Ne Fark\u0131 Var?")
# tablo benzeri iki kolon
panel(s, Inches(0.7), Inches(2.0), Inches(5.85), Inches(4.6), fill=BG_PANEL)
txt(s, Inches(1.0), Inches(2.2), Inches(5.3), Inches(0.5),
    "Genel AI Sohbet Ara\u00e7lar\u0131", size=17, color=MUTED, bold=True)
gen = [
    "Tek model, tek bak\u0131\u015f a\u00e7\u0131s\u0131",
    "\u00c7\u0131kt\u0131 kaynaks\u0131z, do\u011frulanmam\u0131\u015f",
    "Denetim ve izlenebilirlik yok",
    "Maliyet kontrol\u00fc s\u0131n\u0131rl\u0131",
    "Kurumsal y\u00f6neti\u015fim eksik",
    "Her seferinde s\u0131f\u0131rdan ba\u015flar (haf\u0131za yok)",
]
bullets(s, Inches(1.0), Inches(2.8), Inches(5.3), Inches(3.6),
        gen, size=14, color=MUTED, bullet_color=RGBColor(0xFF,0x6B,0x6B), gap=10)

panel(s, Inches(6.75), Inches(2.0), Inches(5.85), Inches(4.6), fill=BG_PANEL,
      line=ACCENT2)
txt(s, Inches(7.05), Inches(2.2), Inches(5.3), Inches(0.5),
    "AgentArmy", size=17, color=ACCENT2, bold=True)
ours = [
    "\u00c7oklu uzman ajan + CEO orkestrasyonu",
    "Kaynakl\u0131, web_search ile do\u011frulanm\u0131\u015f",
    "Tam denetim logu + audit izi",
    "Maliyet defteri ile \u015feffaf b\u00fct\u00e7e",
    "Risk matrisi + onay kuyru\u011fu",
    "Kal\u0131c\u0131 bilgi taban\u0131 (Facts) ile \u00f6\u011frenir",
]
bullets(s, Inches(7.05), Inches(2.8), Inches(5.3), Inches(3.6),
        ours, size=14, color=WHITE, bullet_color=ACCENT2, gap=10)

# =======================================================================
# SLIDE 10 — TEKNOLOJİ
# =======================================================================
s = add_slide()
kicker(s, "Teknoloji")
title(s, "Modern, G\u00fcvenilir Teknoloji Y\u0131\u011f\u0131n\u0131")
tech = [
    ("\u2699", "\u00c7ekirdek Motor", ".NET 8 CLI \u00B7 orkestrasyon, ara\u00e7lar, risk politikas\u0131"),
    ("\uD83E\uDDE0", "Zeka", "OpenAI Responses API \u00B7 GPT-4.1 / GPT-5 \u00B7 web_search"),
    ("\uD83D\uDDA5", "Portal", "React + TypeScript + Vite \u00B7 Tailwind \u00B7 34+ ekran"),
    ("\uD83D\uDDC4", "Veri & Kuyruk", "Supabase (Postgres) \u00B7 RLS \u00B7 31 migration"),
    ("\u2601", "Da\u011f\u0131t\u0131m", "Vercel \u00B7 GitHub Actions worker (\u00fccretsiz \u00e7al\u0131\u015ft\u0131rma)"),
    ("\uD83D\uDD10", "G\u00fcvenlik", "Row-Level Security \u00B7 izin matrisi \u00B7 onay zorlamas\u0131"),
]
cw = Inches(3.85); ch = Inches(1.75)
xs = [Inches(0.7), Inches(4.72), Inches(8.74)]
ys = [Inches(2.0), Inches(3.9)]
idx = 0
for row in ys:
    for col in xs:
        ic, hd, bd = tech[idx]
        card(s, col, row, cw, ch, ic, hd, bd, accent=ACCENT,
             head_size=16, body_size=12.5)
        idx += 1
txt(s, Inches(0.7), Inches(5.95), Inches(12), Inches(0.6),
    "Esnek mimari: \u015firket i\u00e7i (on-prem) veya bulut; kendi modelinizi/anahtar\u0131n\u0131z\u0131 kullanma se\u00e7ene\u011fi.",
    size=14, color=MUTED)

# =======================================================================
# SLIDE 11 — DEĞER / ROI
# =======================================================================
s = add_slide()
kicker(s, "\u0130\u015f De\u011feri", color=ACCENT2)
title(s, "\u00d6l\u00e7\u00fclebilir Kazan\u0131m")
metrics = [
    ("10x", "Daha h\u0131zl\u0131 rapor\n\u00fcretimi", ACCENT),
    ("7/24", "Durmadan \u00e7al\u0131\u015fan\nekip", ACCENT2),
    ("%100", "\u0130zlenebilir &\ndenetlenebilir \u00e7\u0131kt\u0131", ACCENT3),
    ("\u20BA", "\u015eeffaf maliyet,\nkontroll\u00fc b\u00fct\u00e7e", RGBColor(0xB0,0x90,0xFF)),
]
cw = Inches(2.92); ch = Inches(2.2)
xs = [Inches(0.7), Inches(3.78), Inches(6.86), Inches(9.94)]
for col, (big, lbl, cc) in zip(xs, metrics):
    pn = panel(s, col, Inches(2.1), cw, ch)
    txt(s, col, Inches(2.45), cw, Inches(1.0), big, size=44, color=cc,
        bold=True, align=PP_ALIGN.CENTER)
    txt(s, col, Inches(3.55), cw, Inches(0.8), lbl, size=15, color=WHITE,
        align=PP_ALIGN.CENTER, line_spacing=1.05)
txt(s, Inches(0.7), Inches(4.8), Inches(12), Inches(1.6),
    "KPI ile y\u00f6netim: do\u011fruluk, kaynak kapsamas\u0131, tutarl\u0131l\u0131k, h\u0131z, maliyet ve insan d\u00fczeltme oran\u0131.\n"
    "Ekipler tekrar eden bilgi i\u015fini ajanlara devredip strateji ve karara odaklan\u0131r.",
    size=15, color=MUTED, line_spacing=1.2)

# =======================================================================
# SLIDE 12 — KAPANIŞ / CTA
# =======================================================================
s = add_slide()
strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.25), SLIDE_H)
strip.fill.solid(); strip.fill.fore_color.rgb = ACCENT2
strip.line.fill.background(); strip.shadow.inherit = False
txt(s, Inches(0.9), Inches(1.6), Inches(11), Inches(0.5),
    "SONRAKI ADIM", size=15, color=ACCENT2, bold=True)
txt(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.6),
    "Pilot ile Ba\u015flayal\u0131m", size=44, color=WHITE, bold=True)
txt(s, Inches(0.9), Inches(3.6), Inches(11), Inches(1.2),
    "Sizin bir i\u015f ak\u0131\u015f\u0131n\u0131z\u0131 (\u00f6r. haftal\u0131k pazar brief\u2019i veya rakip analizi)\n"
    "2 haftal\u0131k bir pilotta canl\u0131ya alal\u0131m \u2014 \u00e7\u0131kt\u0131y\u0131 birlikte \u00f6l\u00e7elim.",
    size=18, color=MUTED, line_spacing=1.2)
steps = ["1 \u00B7 Ke\u015fif g\u00f6r\u00fc\u015fmesi", "2 \u00B7 Pilot kapsam\u0131", "3 \u00B7 2 hafta canl\u0131", "4 \u00B7 De\u011ferlendirme"]
bx = Inches(0.9)
for st in steps:
    w = Inches(2.85)
    pn = panel(s, bx, Inches(5.0), w, Inches(0.7), fill=BG_PANEL)
    txt(s, bx, Inches(5.18), w, Inches(0.4), st, size=13.5, color=WHITE,
        bold=True, align=PP_ALIGN.CENTER)
    bx = Emu(int(bx) + int(w) + int(Inches(0.12)))
txt(s, Inches(0.9), Inches(6.4), Inches(11), Inches(0.5),
    "AgentArmy \u2014 Sizin yerinize \u00e7al\u0131\u015fan yapay zeka ajan ordusu.", size=15,
    color=ACCENT2, bold=True)

# =======================================================================
# SLIDE 13 — BİZ KİMİZ (Techmora)
# =======================================================================
s = add_slide()
kicker(s, "Biz Kimiz")
title(s, "AgentArmy\u2019nin arkas\u0131nda Techmora var.")

txt(s, Inches(0.7), Inches(1.95), Inches(7.4), Inches(1.3),
    "Techmora; yapay zeka destekli \u00fcr\u00fcnler, siber g\u00fcvenlik ve modern yaz\u0131l\u0131m "
    "m\u00fchendisli\u011fi oda\u011f\u0131nda kurumsal \u00e7\u00f6z\u00fcmler ve finansal altyap\u0131lar geli\u015ftiren "
    "bir teknoloji ve yat\u0131r\u0131m \u015firketidir.",
    size=14.5, color=MUTED, line_spacing=1.18)

txt(s, Inches(0.7), Inches(3.25), Inches(7.4), Inches(0.4),
    "KURUCU ARKA PLANI", size=12, color=ACCENT, bold=True)
# italic vurgu bloğu + sol bar
ib = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(3.62),
                        Inches(0.08), Inches(1.15))
ib.fill.solid(); ib.fill.fore_color.rgb = ACCENT2
ib.line.fill.background(); ib.shadow.inherit = False
txt(s, Inches(0.95), Inches(3.62), Inches(7.15), Inches(1.2),
    "Techmora gen\u00e7 bir \u015firket olsa da; kurucular\u0131n\u0131n ge\u00e7mi\u015fte yer ald\u0131\u011f\u0131 "
    "kurumlar ve y\u00fcr\u00fctt\u00fc\u011f\u00fc projeler, T\u00fcrkiye ve d\u00fcnya \u00e7ap\u0131nda b\u00fcy\u00fck \u00f6l\u00e7ekli "
    "i\u015flerdir.",
    size=14, color=WHITE, italic=True, line_spacing=1.18)

txt(s, Inches(0.7), Inches(4.95), Inches(7.4), Inches(0.4),
    "\u00dcR\u00dcN PORTF\u00d6Y\u00dc", size=12, color=ACCENT, bold=True)
txt(s, Inches(0.7), Inches(5.32), Inches(7.4), Inches(0.8),
    "TERRAM\u0130ND  \u00B7  GUARD\u0130ONX  \u00B7  GFRAUD  \u00B7  AgentArmy  \u00B7  Wallet  \u00B7  "
    "A\u00e7\u0131k Bankac\u0131l\u0131k  \u00B7  Sanal/Fiziki POS  \u00B7  Check\u2019NFit",
    size=14, color=WHITE, line_spacing=1.2)

txt(s, Inches(0.7), Inches(6.5), Inches(7.4), Inches(0.4),
    "\u0130stanbul Teknokent, Avc\u0131lar   \u00B7   info@techmorainvest.com   \u00B7   techmorainvest.com",
    size=12, color=MUTED)

# sağ kartlar
right_cards = [
    ("Mimari Disiplin", "Kurumsal \u00f6l\u00e7ekte s\u00fcrd\u00fcr\u00fclebilirlik; net s\u0131n\u0131rlar, g\u00f6zlemlenebilirlik ve g\u00fcvenli varsay\u0131lanlar.", ACCENT2),
    ("G\u00fcven & Uyum", "KVKK ve sekt\u00f6r standartlar\u0131na uygun, denetlenebilir s\u00fcre\u00e7lerle y\u00f6netilen riskler.", ACCENT),
    ("\u00dcr\u00fcn Odakl\u0131l\u0131k", "\u00d6l\u00e7\u00fcmleme ve iterasyonla de\u011fer \u00fcreten, uzun \u00f6m\u00fcrl\u00fc \u00fcr\u00fcnler.", ACCENT3),
]
ry = Inches(2.0)
for hd, bd, cc in right_cards:
    pn = panel(s, Inches(8.45), ry, Inches(4.15), Inches(1.45))
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.45), ry,
                             Inches(0.1), Inches(1.45))
    bar.fill.solid(); bar.fill.fore_color.rgb = cc
    bar.line.fill.background(); bar.shadow.inherit = False
    txt(s, Inches(8.75), ry + Inches(0.18), Inches(3.7), Inches(0.5), hd,
        size=15.5, color=WHITE, bold=True)
    txt(s, Inches(8.75), ry + Inches(0.62), Inches(3.7), Inches(0.75), bd,
        size=12, color=MUTED, line_spacing=1.08)
    ry = Emu(int(ry) + int(Inches(1.62)))

txt(s, Inches(0.7), Inches(7.0), Inches(8), Inches(0.35),
    "AgentArmy \u2014 M\u00fc\u015fteri Sunumu", size=11, color=MUTED)

# =======================================================================
# SLIDE 14 — KURUCU REFERANSLARI
# =======================================================================
s = add_slide()
kicker(s, "Kurucu Referanslar\u0131")
title(s, "Kurucunun referanslar\u0131")
txt(s, Inches(0.7), Inches(1.9), Inches(11.5), Inches(0.5),
    "Kurucumuzun ge\u00e7mi\u015fte g\u00f6rev ald\u0131\u011f\u0131 kurumlar, \u00e7al\u0131\u015ft\u0131\u011f\u0131 projeler ve i\u015f "
    "birliklerinden bir se\u00e7ki.", size=14, color=MUTED)


def chip_row(slide, top, label, items, color):
    txt(slide, Inches(0.7), top, Inches(11.5), Inches(0.35),
        label, size=12, color=color, bold=True)
    cx = Inches(0.7)
    cy = Emu(int(top) + int(Inches(0.4)))
    for it in items:
        w = Inches(0.42 + 0.105 * len(it))
        if int(cx) + int(w) > int(Inches(12.7)):
            cx = Inches(0.7)
            cy = Emu(int(cy) + int(Inches(0.62)))
        ch = panel(slide, cx, cy, w, Inches(0.5), fill=BG_PANEL, line=color)
        txt(slide, cx, cy + Inches(0.11), w, Inches(0.3), it, size=12.5,
            color=WHITE, align=PP_ALIGN.CENTER)
        cx = Emu(int(cx) + int(w) + int(Inches(0.18)))
    return Emu(int(cy) + int(Inches(0.62)))


ny = Inches(2.65)
ny = chip_row(s, ny, "KAMU & SAVUNMA",
              ["\u00c7evre ve \u015eehircilik Bakanl\u0131\u011f\u0131", "\u0130\u00e7i\u015fleri Bakanl\u0131\u011f\u0131",
               "T\u00dc\u0130K", "Yarg\u0131tay", "Aselsan"], ACCENT2)
ny = Emu(int(ny) + int(Inches(0.18)))
ny = chip_row(s, ny, "\u00d6ZEL SEKT\u00d6R & F\u0130NANS",
              ["\u0130\u015fbir Holding", "Bitexen", "Dinamik Yat\u0131r\u0131m", "Arsa ve Ev",
               "Eslior", "Fast Company"], ACCENT)
ny = Emu(int(ny) + int(Inches(0.18)))
ny = chip_row(s, ny, "ULUSLARARASI", ["Nikon", "Canon"], ACCENT3)
ny = Emu(int(ny) + int(Inches(0.18)))
ny = chip_row(s, ny, "\u00dcN\u0130VERS\u0130TELER",
              ["Orta Do\u011fu Teknik \u00dcniversitesi", "K\u0131r\u0131kkale \u00dcniversitesi",
               "Ege \u00dcniversitesi", "Dokuz Eyl\u00fcl \u00dcniversitesi"],
              RGBColor(0xB0, 0x90, 0xFF))

txt(s, Inches(0.7), Emu(int(ny) + int(Inches(0.05))), Inches(11.5), Inches(0.4),
    "+ ve \u00e7e\u015fitli k\u00fc\u00e7\u00fck \u00e7apl\u0131 e-ticaret sistemleri", size=12.5,
    color=MUTED, italic=True)

# ---- kaydet
import os
out = os.path.join(os.path.dirname(__file__), "AgentArmy-Satis-Sunumu.pptx")
prs.save(out)
print("KAYDED\u0130LD\u0130:", out)
print("Slayt say\u0131s\u0131:", len(prs.slides._sldIdLst))
